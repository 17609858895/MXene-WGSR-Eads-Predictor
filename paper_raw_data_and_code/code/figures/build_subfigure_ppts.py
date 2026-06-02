from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.28
GAP = 0.16


def contain(image_path: Path, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    with Image.open(image_path) as im:
        img_ratio = im.width / im.height
    box_ratio = w / h
    if img_ratio > box_ratio:
        out_w = w
        out_h = w / img_ratio
        return x, y + (h - out_h) / 2, out_w, out_h
    out_h = h
    out_w = h * img_ratio
    return x + (w - out_w) / 2, y, out_w, out_h


def add_image(slide, folder: Path, file_name: str, x: float, y: float, w: float, h: float) -> None:
    image_path = folder / file_name
    if not image_path.exists():
        raise FileNotFoundError(image_path)
    px, py, pw, ph = contain(image_path, x, y, w, h)
    slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))


def make_deck(fig_no: int, boxes: list[dict[str, float | str]]) -> Path:
    folder = ROOT / f"Fig{fig_no:02d}"
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    for box in boxes:
        add_image(slide, folder, str(box["file"]), float(box["x"]), float(box["y"]), float(box["w"]), float(box["h"]))
    out = folder / f"Fig{fig_no:02d}_subfigures_layout.pptx"
    prs.save(out)
    return out


def grid(
    files: list[str],
    rows: int,
    cols: int,
    *,
    left: float = MARGIN,
    top: float = MARGIN,
    right: float = MARGIN,
    bottom: float = MARGIN,
    gap_x: float = GAP,
    gap_y: float = GAP,
) -> list[dict[str, float | str]]:
    w = (SLIDE_W - left - right - gap_x * (cols - 1)) / cols
    h = (SLIDE_H - top - bottom - gap_y * (rows - 1)) / rows
    return [
        {
            "file": file,
            "x": left + (i % cols) * (w + gap_x),
            "y": top + (i // cols) * (h + gap_y),
            "w": w,
            "h": h,
        }
        for i, file in enumerate(files)
    ]


def main() -> None:
    outputs: list[Path] = []

    outputs.append(
        make_deck(
            2,
            grid(
                [
                    "Fig02_panel_a_label_availability.png",
                    "Fig02_panel_b_eads_histogram.png",
                    "Fig02_panel_c_eads_by_adsorbate.png",
                    "Fig02_panel_d_pearson_correlation.png",
                    "Fig02_panel_e_spearman_correlation.png",
                    "Fig02_panel_f_mutual_information.png",
                ],
                2,
                3,
            ),
        )
    )

    outputs.append(
        make_deck(
            3,
            grid(
                [
                    "Fig03_panel_a_permutation_importance.png",
                    "Fig03_panel_b_feature_dendrogram.png",
                    "Fig03_panel_c_consensus_diagram.png",
                ],
                1,
                3,
                top=0.45,
                bottom=0.45,
                gap_x=0.20,
            ),
        )
    )

    outputs.append(
        make_deck(
            4,
            [
                *grid(
                    [
                        "Fig04_panel_a_rf.png",
                        "Fig04_panel_b_extratrees.png",
                        "Fig04_panel_c_gbr.png",
                        "Fig04_panel_d_hgbr.png",
                        "Fig04_panel_e_svr.png",
                        "Fig04_panel_f_mlp.png",
                    ],
                    2,
                    3,
                    top=0.20,
                    bottom=1.72,
                    gap_x=0.10,
                    gap_y=0.12,
                ),
                {"file": "Fig04_panel_g_metrics_comparison.png", "x": 0.42, "y": 5.94, "w": 8.35, "h": 1.35},
                {"file": "Fig04_panel_h_learning_curve.png", "x": 9.15, "y": 5.94, "w": 3.68, "h": 1.35},
            ],
        )
    )

    outputs.append(
        make_deck(
            5,
            grid(
                [
                    "Fig05_panel_a_topsis_weights.png",
                    "Fig05_panel_b_error_distribution.png",
                    "Fig05_panel_c_wilcoxon_heatmap.png",
                ],
                1,
                3,
                top=0.45,
                bottom=0.45,
                gap_x=0.22,
            ),
        )
    )

    outputs.append(
        make_deck(
            6,
            [
                {"file": "Fig06_panel_a_contribution_scatter.png", "x": 0.28, "y": 0.30, "w": 4.15, "h": 6.82},
                {"file": "Fig06_panel_b_permutation_importance.png", "x": 4.62, "y": 0.30, "w": 4.05, "h": 3.25},
                {"file": "Fig06_panel_c_pdp.png", "x": 8.95, "y": 0.30, "w": 4.05, "h": 3.25},
                {"file": "Fig06_panel_d_ale.png", "x": 4.62, "y": 3.86, "w": 4.05, "h": 3.25},
                {"file": "Fig06_panel_e_interaction_dependency.png", "x": 8.95, "y": 3.86, "w": 4.05, "h": 3.25},
            ],
        )
    )

    outputs.append(
        make_deck(
            7,
            grid(
                [
                    "Fig07_panel_a_causal_dag.png",
                    "Fig07_panel_b_counterfactual_effects.png",
                    "Fig07_panel_c_ads_vs_barrier.png",
                ],
                1,
                3,
                top=0.45,
                bottom=0.45,
                gap_x=0.20,
            ),
        )
    )

    outputs.append(
        make_deck(
            8,
            [
                *grid(
                    [
                        "Fig08_panel_a_transfer_scatter.png",
                        "Fig08_panel_b_residual_boxplot.png",
                        "Fig08_panel_c_dft_score_barh.png",
                    ],
                    1,
                    3,
                    top=0.25,
                    bottom=4.28,
                    gap_x=0.20,
                ),
                {"file": "Fig08_panel_d_periodic_heatmap.png", "x": 0.70, "y": 3.42, "w": 11.95, "h": 3.68},
            ],
        )
    )

    for output in outputs:
        print(output)


if __name__ == "__main__":
    main()
